# Copyright 2025 Google LLC
# SPDX-License-Identifier: Apache-2.0

"""PowerPoint presentation generator service.

Creates .pptx files for reports, proposals, and pitch decks
using python-pptx library.
"""

import os
from datetime import datetime
from typing import Any
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


class PptxGenerator:
    """Service for generating PowerPoint presentations.
    
    Supports:
    - Title slides
    - Content slides with bullet points
    - Data tables
    - Summary slides
    - Custom branding
    """
    
    # Default colors (can be customized per user/company)
    PRIMARY_COLOR = RGBColor(0, 128, 128)  # Teal
    SECONDARY_COLOR = RGBColor(64, 64, 64)  # Dark gray
    ACCENT_COLOR = RGBColor(0, 176, 176)  # Light teal
    
    def __init__(self, output_dir: str | None = None):
        """Initialize the generator.
        
        Args:
            output_dir: Directory for saving generated files.
                       Defaults to .tmp/reports/
        """
        self.output_dir = output_dir or os.path.join(
            os.path.dirname(os.path.dirname(os.path.dirname(__file__))),
            ".tmp",
            "reports",
        )
        Path(self.output_dir).mkdir(parents=True, exist_ok=True)
    
    def create_report_presentation(
        self,
        title: str,
        subtitle: str = "",
        sections: list[dict[str, Any]] | None = None,
        author: str = "Pikar AI",
    ) -> str:
        """Create a report presentation.
        
        Args:
            title: Main title for the presentation.
            subtitle: Optional subtitle.
            sections: List of section dicts with keys:
                - title: Section/slide title
                - type: "bullets" | "table" | "summary"
                - content: List of bullet points or table data
            author: Author name for title slide.
            
        Returns:
            Path to the generated .pptx file.
        """
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
        prs.slide_height = Inches(7.5)
        
        # Title slide
        self._add_title_slide(prs, title, subtitle, author)
        
        # Content slides
        if sections:
            for section in sections:
                section_type = section.get("type", "bullets")
                if section_type == "bullets":
                    self._add_bullet_slide(
                        prs,
                        section.get("title", ""),
                        section.get("content", []),
                    )
                elif section_type == "table":
                    self._add_table_slide(
                        prs,
                        section.get("title", ""),
                        section.get("headers", []),
                        section.get("rows", []),
                    )
                elif section_type == "summary":
                    self._add_summary_slide(
                        prs,
                        section.get("title", "Key Takeaways"),
                        section.get("content", []),
                    )
        
        # Generate filename
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_title = "".join(c if c.isalnum() else "_" for c in title)[:30]
        filename = f"{safe_title}_{timestamp}.pptx"
        filepath = os.path.join(self.output_dir, filename)
        
        prs.save(filepath)
        return filepath
    
    def create_daily_report(
        self,
        date: str,
        metrics: dict[str, Any],
        highlights: list[str],
        concerns: list[str] | None = None,
        recommendations: list[str] | None = None,
    ) -> str:
        """Create a daily report presentation.
        
        Args:
            date: Report date string.
            metrics: Dict of metric name -> value.
            highlights: List of positive highlights.
            concerns: Optional list of concerns or issues.
            recommendations: Optional list of action items.
            
        Returns:
            Path to the generated .pptx file.
        """
        sections = [
            {
                "title": "Key Metrics",
                "type": "table",
                "headers": ["Metric", "Value"],
                "rows": [[k, str(v)] for k, v in metrics.items()],
            },
            {
                "title": "Highlights",
                "type": "bullets",
                "content": highlights,
            },
        ]
        
        if concerns:
            sections.append({
                "title": "Concerns & Issues",
                "type": "bullets",
                "content": concerns,
            })
        
        if recommendations:
            sections.append({
                "title": "Recommendations",
                "type": "summary",
                "content": recommendations,
            })
        
        return self.create_report_presentation(
            title=f"Daily Report - {date}",
            subtitle="End of Day Summary",
            sections=sections,
        )
    
    def create_weekly_report(
        self,
        week_range: str,
        summary: str,
        metrics_comparison: list[dict[str, Any]],
        achievements: list[str],
        next_week_priorities: list[str],
    ) -> str:
        """Create a weekly report presentation.
        
        Args:
            week_range: Week range string (e.g., "Jan 20 - Jan 26, 2026").
            summary: Executive summary text.
            metrics_comparison: List of {metric, previous, current, change}.
            achievements: List of week's achievements.
            next_week_priorities: List of priorities for next week.
            
        Returns:
            Path to the generated .pptx file.
        """
        sections = [
            {
                "title": "Executive Summary",
                "type": "bullets",
                "content": [summary],
            },
            {
                "title": "Metrics Comparison",
                "type": "table",
                "headers": ["Metric", "Previous", "Current", "Change"],
                "rows": [
                    [m["metric"], str(m["previous"]), str(m["current"]), m["change"]]
                    for m in metrics_comparison
                ],
            },
            {
                "title": "Achievements",
                "type": "bullets",
                "content": achievements,
            },
            {
                "title": "Next Week Priorities",
                "type": "summary",
                "content": next_week_priorities,
            },
        ]
        
        return self.create_report_presentation(
            title=f"Weekly Report",
            subtitle=week_range,
            sections=sections,
        )
    
    def _add_title_slide(
        self,
        prs: Presentation,
        title: str,
        subtitle: str,
        author: str,
    ) -> None:
        """Add a title slide."""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.paragraphs[0].text = title
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.PRIMARY_COLOR
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Subtitle
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4), Inches(12.333), Inches(0.75)
            )
            sub_frame = sub_box.text_frame
            sub_frame.paragraphs[0].text = subtitle
            sub_frame.paragraphs[0].font.size = Pt(24)
            sub_frame.paragraphs[0].font.color.rgb = self.SECONDARY_COLOR
            sub_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Author and date
        footer_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5)
        )
        footer_frame = footer_box.text_frame
        date_str = datetime.now().strftime("%B %d, %Y")
        footer_frame.paragraphs[0].text = f"{author} • {date_str}"
        footer_frame.paragraphs[0].font.size = Pt(14)
        footer_frame.paragraphs[0].font.color.rgb = self.SECONDARY_COLOR
        footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _add_bullet_slide(
        self,
        prs: Presentation,
        title: str,
        bullets: list[str],
    ) -> None:
        """Add a slide with bullet points."""
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12.333), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.paragraphs[0].text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.PRIMARY_COLOR
        
        # Bullets
        content_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(1.75), Inches(11.583), Inches(5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(20)
            p.font.color.rgb = self.SECONDARY_COLOR
            p.space_after = Pt(12)
    
    def _add_table_slide(
        self,
        prs: Presentation,
        title: str,
        headers: list[str],
        rows: list[list[str]],
    ) -> None:
        """Add a slide with a data table."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12.333), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.paragraphs[0].text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.PRIMARY_COLOR
        
        # Table
        if not headers or not rows:
            return
        
        col_count = len(headers)
        row_count = len(rows) + 1  # +1 for header
        
        table = slide.shapes.add_table(
            row_count, col_count,
            Inches(0.75), Inches(1.75),
            Inches(11.583), Inches(min(row_count * 0.5, 5))
        ).table
        
        # Header row
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(14)
        
        # Data rows
        for row_idx, row in enumerate(rows):
            for col_idx, value in enumerate(row):
                if col_idx < col_count:
                    cell = table.cell(row_idx + 1, col_idx)
                    cell.text = str(value)
                    cell.text_frame.paragraphs[0].font.size = Pt(12)
    
    def _add_summary_slide(
        self,
        prs: Presentation,
        title: str,
        points: list[str],
    ) -> None:
        """Add a summary/key takeaways slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Title with icon-like marker
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12.333), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.paragraphs[0].text = f"🎯 {title}"
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.ACCENT_COLOR
        
        # Key points with numbers
        content_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(1.75), Inches(11.583), Inches(5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for i, point in enumerate(points):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            p.text = f"{i + 1}. {point}"
            p.font.size = Pt(20)
            p.font.color.rgb = self.SECONDARY_COLOR
            p.space_after = Pt(16)


# Singleton instance
pptx_generator = PptxGenerator()


def generate_report_pptx(
    title: str,
    sections: list[dict[str, Any]],
    subtitle: str = "",
) -> str:
    """Generate a PowerPoint report.
    
    Args:
        title: Report title.
        sections: List of section configurations.
        subtitle: Optional subtitle.
        
    Returns:
        Path to the generated .pptx file.
    """
    return pptx_generator.create_report_presentation(title, subtitle, sections)
