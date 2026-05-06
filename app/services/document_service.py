# Copyright (c) 2024-2026 Pikar AI. All rights reserved.
# Proprietary and confidential. See LICENSE file for details.

"""Document generation service — PDF reports, PowerPoint pitch decks, and spreadsheets.

Generates branded PDF documents from Jinja2 HTML templates using weasyprint,
and branded PowerPoint presentations using python-pptx.  All heavy rendering
runs in a thread-pool to avoid blocking the async event loop.

Documents are uploaded to Supabase Storage and tracked in the ``media_assets``
table for conversation linking.
"""

from __future__ import annotations

import asyncio
import base64
import html
import logging
import uuid
from io import BytesIO
from pathlib import Path
from typing import Any

from jinja2 import Environment, FileSystemLoader
from markupsafe import Markup

from app.agents.tools.brand_profile import get_brand_profile
from app.rag.knowledge_vault import ingest_document_content
from app.services.document_text_extraction import (
    ExtractionError,
    extract_text_from_bytes,
)
from app.services.supabase import get_service_client
from app.services.supabase_async import execute_async

logger = logging.getLogger(__name__)


def _get_matplotlib():
    """Lazy import matplotlib to avoid import-time side effects."""
    import matplotlib

    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    return plt


def _get_weasyprint_html():
    """Lazy import weasyprint.HTML (requires system C libraries)."""
    from weasyprint import HTML as _HTML

    return _HTML

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

VALID_TEMPLATES = [
    "financial_report",
    "project_proposal",
    "meeting_summary",
    "competitive_analysis",
    "sales_proposal",
    "narrative_report",
]

TEMPLATE_DIR = Path(__file__).resolve().parent.parent / "templates" / "pdf"
DOCUMENT_BUCKET = "generated-documents"

# Pikar default brand colors (used when user has no brand profile)
DEFAULT_BRAND_NAME = "Pikar"
DEFAULT_PRIMARY_COLOR = "#4F46E5"
DEFAULT_SECONDARY_COLOR = "#10B981"
DEFAULT_ACCENT_COLOR = "#F59E0B"

# Safety limit: reject PDFs larger than 5 MB (~50+ pages of content)
MAX_PDF_BYTES = 5 * 1024 * 1024


# ---------------------------------------------------------------------------
# Service
# ---------------------------------------------------------------------------


class DocumentService:
    """Generate branded PDF, PPTX, and XLSX documents from structured data."""

    def __init__(self) -> None:
        self._jinja_env = Environment(
            loader=FileSystemLoader(str(TEMPLATE_DIR)),
            autoescape=True,
        )
        self._jinja_env.filters["markdown"] = self._render_markdown

    @staticmethod
    def _render_markdown(text: str | None) -> Markup:
        """Render markdown text to safe HTML for embedding in PDF templates.

        Raw HTML in the source is escaped (``html=False``), so agent-provided
        content cannot inject arbitrary tags into the PDF.
        """
        if not text:
            return Markup("")
        from markdown_it import MarkdownIt

        md = MarkdownIt("commonmark", {"html": False, "linkify": True, "breaks": False})
        md.enable(["table", "strikethrough"])
        return Markup(md.render(str(text)))

    # ------------------------------------------------------------------
    # PDF Generation
    # ------------------------------------------------------------------

    async def generate_pdf(
        self,
        template_name: str,
        data: dict[str, Any],
        user_id: str,
        session_id: str | None = None,
        title: str | None = None,
    ) -> dict[str, Any]:
        """Render an HTML template to a branded PDF and upload it.

        Args:
            template_name: One of :data:`VALID_TEMPLATES`.
            data: Template-specific data dict.
            user_id: Owner of the generated document.
            session_id: Optional conversation session for linking.
            title: Human-readable document title.

        Returns:
            Widget dict suitable for returning to the chat UI.

        Raises:
            ValueError: If *template_name* is invalid or the PDF exceeds size limits.
        """
        if template_name not in VALID_TEMPLATES:
            msg = f"Invalid template '{template_name}'. Must be one of {VALID_TEMPLATES}"
            raise ValueError(msg)

        # Resolve branding ---------------------------------------------------
        brand = await get_brand_profile(user_id=user_id)
        brand_vars = self._extract_brand_vars(brand)

        # Render charts if present -------------------------------------------
        chart_images: dict[str, str] = {}
        chart_entries = data.get("chart_data", [])
        if isinstance(chart_entries, list):
            for idx, chart_data in enumerate(chart_entries):
                png_bytes = self.render_chart(chart_data)
                b64 = base64.b64encode(png_bytes).decode()
                chart_images[f"chart_{idx}"] = f"data:image/png;base64,{b64}"

        # Build template context ---------------------------------------------
        context: dict[str, Any] = {
            **data,
            **brand_vars,
            "chart_images": chart_images,
        }

        template = self._jinja_env.get_template(f"{template_name}.html")
        rendered_html = template.render(**context)

        # Render PDF in thread pool (weasyprint is CPU-intensive sync code) --
        HTMLClass = _get_weasyprint_html()
        pdf_bytes: bytes = await asyncio.to_thread(
            HTMLClass(string=rendered_html).write_pdf,
        )

        # Enforce page/size limit --------------------------------------------
        if len(pdf_bytes) > MAX_PDF_BYTES:
            msg = (
                f"Generated PDF ({len(pdf_bytes):,} bytes) exceeds maximum "
                f"allowed size ({MAX_PDF_BYTES:,} bytes / ~50 pages)."
            )
            raise ValueError(msg)

        # Upload + track -----------------------------------------------------
        doc_id = str(uuid.uuid4())
        doc_title = title or f"{template_name.replace('_', ' ').title()}"
        filename = f"{doc_id}.pdf"

        widget = await self._upload_document(
            file_bytes=pdf_bytes,
            user_id=user_id,
            doc_id=doc_id,
            filename=filename,
            content_type="application/pdf",
            title=doc_title,
            template_name=template_name,
            session_id=session_id,
            file_type="pdf",
        )
        return widget

    # ------------------------------------------------------------------
    # PPTX Generation
    # ------------------------------------------------------------------

    async def generate_pptx(
        self,
        slides_data: list[dict[str, Any]],
        user_id: str,
        session_id: str | None = None,
        title: str | None = None,
    ) -> dict[str, Any]:
        """Generate a branded PowerPoint pitch deck.

        Args:
            slides_data: List of slide dicts, each with ``title``, ``bullets``,
                and optionally ``chart_image_bytes``.
            user_id: Owner of the generated document.
            session_id: Optional conversation session for linking.
            title: Human-readable document title.

        Returns:
            Widget dict suitable for returning to the chat UI.
        """
        brand = await get_brand_profile(user_id=user_id)
        brand_vars = self._extract_brand_vars(brand)

        # Build presentation in thread pool ----------------------------------
        pptx_bytes = await asyncio.to_thread(
            self._build_pptx, slides_data, brand_vars,
        )

        doc_id = str(uuid.uuid4())
        doc_title = title or "Pitch Deck"
        filename = f"{doc_id}.pptx"

        widget = await self._upload_document(
            file_bytes=pptx_bytes,
            user_id=user_id,
            doc_id=doc_id,
            filename=filename,
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            title=doc_title,
            template_name="pitch_deck",
            session_id=session_id,
            file_type="pptx",
        )
        return widget

    async def generate_xlsx(
        self,
        sheets_data: list[dict[str, Any]],
        user_id: str,
        session_id: str | None = None,
        title: str | None = None,
    ) -> dict[str, Any]:
        """Generate a branded Excel workbook and upload it for download."""
        brand = await get_brand_profile(user_id=user_id)
        brand_vars = self._extract_brand_vars(brand)

        workbook_bytes = await asyncio.to_thread(
            self._build_xlsx,
            sheets_data,
            brand_vars,
        )

        doc_id = str(uuid.uuid4())
        doc_title = title or "Spreadsheet Export"
        filename = f"{doc_id}.xlsx"

        widget = await self._upload_document(
            file_bytes=workbook_bytes,
            user_id=user_id,
            doc_id=doc_id,
            filename=filename,
            content_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            title=doc_title,
            template_name="spreadsheet_export",
            session_id=session_id,
            file_type="xlsx",
        )
        return widget

    @staticmethod
    def _build_pptx(
        slides_data: list[dict[str, Any]],
        brand_vars: dict[str, Any],
    ) -> bytes:
        """Build a PPTX file in memory (sync, runs in thread pool)."""
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt

        prs = Presentation()
        # 16:9 widescreen
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        primary_hex = brand_vars.get("primary_color", DEFAULT_PRIMARY_COLOR).lstrip("#")
        primary_rgb = RGBColor(
            int(primary_hex[0:2], 16),
            int(primary_hex[2:4], 16),
            int(primary_hex[4:6], 16),
        )

        blank_layout = prs.slide_layouts[6]  # blank layout

        for slide_data in slides_data:
            slide = prs.slides.add_slide(blank_layout)

            # Title textbox
            left = Inches(0.8)
            top = Inches(0.5)
            width = Inches(11.5)
            height = Inches(1.2)
            txbox = slide.shapes.add_textbox(left, top, width, height)
            tf = txbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data.get("title", "")
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = primary_rgb

            # Bullet points
            bullets = slide_data.get("bullets", [])
            if bullets:
                bullet_top = Inches(2.0)
                bullet_height = Inches(4.0)
                btxbox = slide.shapes.add_textbox(left, bullet_top, width, bullet_height)
                btf = btxbox.text_frame
                btf.word_wrap = True
                for i, bullet in enumerate(bullets):
                    if i == 0:
                        bp = btf.paragraphs[0]
                    else:
                        bp = btf.add_paragraph()
                    bp.text = bullet
                    bp.font.size = Pt(20)
                    bp.space_after = Pt(12)

            # Chart image (if provided)
            chart_bytes = slide_data.get("chart_image_bytes")
            if chart_bytes and isinstance(chart_bytes, bytes):
                try:
                    img_stream = BytesIO(chart_bytes)
                    img_left = Inches(8.0)
                    img_top = Inches(2.0)
                    img_width = Inches(4.5)
                    slide.shapes.add_picture(img_stream, img_left, img_top, width=img_width)
                except Exception:
                    logger.warning("Failed to embed chart image in PPTX slide")

        buf = BytesIO()
        prs.save(buf)
        return buf.getvalue()

    @staticmethod
    def _build_xlsx(
        sheets_data: list[dict[str, Any]],
        brand_vars: dict[str, Any],
    ) -> bytes:
        """Build an XLSX workbook in memory (sync, runs in thread pool)."""
        from openpyxl import Workbook
        from openpyxl.styles import Alignment, Font, PatternFill
        from openpyxl.utils import get_column_letter

        workbook = Workbook()
        default_sheet = workbook.active
        workbook.remove(default_sheet)

        primary_hex = brand_vars.get("primary_color", DEFAULT_PRIMARY_COLOR).lstrip("#")
        header_fill = PatternFill(fill_type="solid", fgColor=primary_hex)
        header_font = Font(color="FFFFFF", bold=True)
        wrap_alignment = Alignment(vertical="top", wrap_text=True)

        normalized_sheets = sheets_data or [
            {
                "name": "Sheet1",
                "headers": ["Value"],
                "rows": [["No data provided"]],
            }
        ]

        for index, sheet_data in enumerate(normalized_sheets, start=1):
            raw_name = str(sheet_data.get("name") or f"Sheet{index}").strip() or f"Sheet{index}"
            worksheet = workbook.create_sheet(title=raw_name[:31])

            headers = sheet_data.get("headers") or []
            rows = sheet_data.get("rows") or []
            title = str(sheet_data.get("title") or "").strip()

            current_row = 1
            if title:
                worksheet.cell(row=current_row, column=1, value=title)
                worksheet.cell(row=current_row, column=1).font = Font(bold=True, size=14)
                current_row += 2

            if headers:
                for column_index, header in enumerate(headers, start=1):
                    cell = worksheet.cell(row=current_row, column=column_index, value=str(header))
                    cell.fill = header_fill
                    cell.font = header_font
                    cell.alignment = wrap_alignment
                worksheet.freeze_panes = worksheet.cell(row=current_row + 1, column=1)
                worksheet.auto_filter.ref = (
                    f"A{current_row}:"
                    f"{get_column_letter(max(len(headers), 1))}{current_row}"
                )
                current_row += 1

            for row in rows:
                normalized_row = row if isinstance(row, list) else [row]
                for column_index, value in enumerate(normalized_row, start=1):
                    cell = worksheet.cell(row=current_row, column=column_index, value=value)
                    cell.alignment = wrap_alignment
                current_row += 1

            max_columns = max(
                len(headers),
                max(
                    ((len(row) if isinstance(row, list) else 1) for row in rows),
                    default=0,
                ),
                1,
            )
            for column_index in range(1, max_columns + 1):
                values = []
                for row_cells in worksheet.iter_rows(
                    min_col=column_index,
                    max_col=column_index,
                    values_only=True,
                ):
                    cell_value = row_cells[0]
                    if cell_value is not None:
                        values.append(str(cell_value))
                width = min(max((len(value) for value in values), default=10) + 2, 42)
                worksheet.column_dimensions[get_column_letter(column_index)].width = width

        buf = BytesIO()
        workbook.save(buf)
        return buf.getvalue()

    # ------------------------------------------------------------------
    # Chart Rendering
    # ------------------------------------------------------------------

    @staticmethod
    def render_chart(chart_data: dict[str, Any]) -> bytes:
        """Render a chart to PNG bytes using matplotlib.

        Args:
            chart_data: Dict with keys ``type`` (bar|line|pie), ``labels``,
                ``values``, and ``title``.

        Returns:
            PNG image bytes.
        """
        plt = _get_matplotlib()

        chart_type = chart_data.get("type", "bar")
        labels = chart_data.get("labels", [])
        values = chart_data.get("values", [])
        chart_title = chart_data.get("title", "")

        fig, ax = plt.subplots(figsize=(8, 5))
        try:
            if chart_type == "pie":
                ax.pie(values, labels=labels, autopct="%1.1f%%", startangle=90)
            elif chart_type == "line":
                ax.plot(labels, values, marker="o", linewidth=2)
                ax.set_ylabel("Value")
            else:
                # Default to bar chart
                ax.bar(labels, values)
                ax.set_ylabel("Value")

            ax.set_title(chart_title)
            fig.tight_layout()

            buf = BytesIO()
            fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
            buf.seek(0)
            return buf.getvalue()
        finally:
            plt.close(fig)

    # ------------------------------------------------------------------
    # Brand Extraction
    # ------------------------------------------------------------------

    @staticmethod
    def _extract_brand_vars(brand_result: dict[str, Any]) -> dict[str, Any]:
        """Extract brand variables from a get_brand_profile result.

        Falls back to Pikar defaults when no profile is available.
        """
        profile = brand_result.get("profile") if brand_result.get("success") else None

        if profile is None:
            return {
                "brand_name": DEFAULT_BRAND_NAME,
                "logo_url": "",
                "primary_color": DEFAULT_PRIMARY_COLOR,
                "secondary_color": DEFAULT_SECONDARY_COLOR,
                "accent_color": DEFAULT_ACCENT_COLOR,
            }

        visual = profile.get("visual_style") or {}
        palette = visual.get("color_palette") or []

        return {
            "brand_name": profile.get("brand_name", DEFAULT_BRAND_NAME),
            "logo_url": profile.get("logo_url", ""),
            "primary_color": palette[0] if len(palette) > 0 else DEFAULT_PRIMARY_COLOR,
            "secondary_color": palette[1] if len(palette) > 1 else DEFAULT_SECONDARY_COLOR,
            "accent_color": palette[2] if len(palette) > 2 else DEFAULT_ACCENT_COLOR,
        }

    # ------------------------------------------------------------------
    # Upload + Track
    # ------------------------------------------------------------------

    async def _upload_document(
        self,
        file_bytes: bytes,
        user_id: str,
        doc_id: str,
        filename: str,
        content_type: str,
        title: str,
        template_name: str,
        session_id: str | None,
        file_type: str,
    ) -> dict[str, Any]:
        """Upload document to Supabase Storage and track in media_assets.

        Follows the same pattern as director_service for consistency and
        best-effort ingests generated documents into the Knowledge Vault.

        Returns:
            Widget dict for the chat UI.
        """
        supabase = get_service_client()
        path = f"{user_id}/{filename}"

        # Upload to Storage ---------------------------------------------------
        await asyncio.to_thread(
            supabase.storage.from_(DOCUMENT_BUCKET).upload,
            path,
            file_bytes,
            {"content-type": content_type},
        )

        # Signed URL for private access ---------------------------------------
        signed = await asyncio.to_thread(
            supabase.storage.from_(DOCUMENT_BUCKET).create_signed_url,
            path,
            3600,  # 1 hour expiry
        )
        signed_url = signed.get("signedURL", signed.get("signedUrl", ""))

        # Track in media_assets -----------------------------------------------
        try:
            await execute_async(
                supabase.table("media_assets").upsert(
                    {
                        "id": doc_id,
                        "user_id": user_id,
                        "bucket_id": DOCUMENT_BUCKET,
                        "asset_type": "document",
                        "title": title,
                        "filename": filename,
                        "file_path": path,
                        "file_url": signed_url,
                        "file_type": content_type,
                        "category": "generated",
                        "size_bytes": len(file_bytes),
                        "metadata": {
                            "template": template_name,
                            "session_id": session_id,
                        },
                    },
                    on_conflict="id",
                ),
                op_name="document_service.media_assets.upsert",
            )
        except Exception as exc:
            logger.warning("Failed to track document in media_assets: %s", exc)

        # Knowledge Vault auto-ingest (best-effort, non-blocking) -------------
        try:
            if template_name == "pitch_deck":
                document_type = "pitch_deck"
            elif file_type == "xlsx":
                document_type = "spreadsheet"
            else:
                document_type = "pdf"

            ingest_content: str | None = None
            if content_type in {
                "application/pdf",
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            }:
                try:
                    extracted = extract_text_from_bytes(
                        file_bytes,
                        content_type,
                        filename=filename,
                    )
                    if extracted:
                        ingest_content = extracted
                except ExtractionError as exc:
                    logger.warning(
                        "document text extraction failed for %s: %s",
                        doc_id,
                        exc,
                    )

            if not ingest_content:
                if document_type == "pitch_deck":
                    ingest_content = f"Generated pitch deck: {title}. Asset ID: {doc_id}."
                elif document_type == "spreadsheet":
                    ingest_content = (
                        f"Generated spreadsheet export: {title}. "
                        f"Asset ID: {doc_id}."
                    )
                else:
                    ingest_content = (
                        f"Generated PDF report ({template_name}): {title}. "
                        f"Asset ID: {doc_id}."
                    )

            ingest_metadata = {
                "asset_id": doc_id,
                "asset_type": "document",
                "bucket_id": DOCUMENT_BUCKET,
                "file_path": path,
                "template": template_name,
                "file_type": file_type,
                "session_id": session_id,
            }

            await ingest_document_content(
                content=ingest_content,
                title=title,
                document_type=document_type,
                user_id=user_id,
                metadata=ingest_metadata,
            )
        except Exception as exc:
            logger.warning("Knowledge vault ingest for %s failed: %s", file_type, exc)

        # Return widget dict --------------------------------------------------
        return {
            "type": "document",
            "title": title,
            "data": {
                "documentUrl": signed_url,
                "title": title,
                "fileType": file_type,
                "sizeBytes": len(file_bytes),
                "templateName": template_name,
            },
            "dismissible": True,
            "expandable": False,
        }


# ---------------------------------------------------------------------------
# Pure render helpers — used by the document_editor edit tools to re-render
# documents from a source dict after each agent mutation. These functions are
# intentionally side-effect free: they return raw bytes without uploading or
# tracking anywhere. The 4 source schemas are stable v1 contracts shared with
# the document_source/document_version services.
# ---------------------------------------------------------------------------


async def render_pdf_from_source(source: dict[str, Any]) -> bytes:
    """Render a markdown-style source dict to PDF bytes.

    Args:
        source: Dict shaped ``{"title": str, "sections": [{"heading": str,
            "content": str}]}``. ``content`` is split on blank lines into
            paragraphs. All user content is HTML-escaped before rendering.

    Returns:
        PDF file bytes.

    Raises:
        ValueError: If ``source`` is missing the required ``"sections"`` key.
    """
    if "sections" not in source:
        msg = "render_pdf_from_source: source is missing required key 'sections'"
        raise ValueError(msg)

    title = html.escape(str(source.get("title", "Untitled")))
    section_html_parts: list[str] = []
    for section in source["sections"]:
        heading = html.escape(str(section.get("heading", "")))
        content = str(section.get("content", ""))
        paragraph_html = "".join(
            f"<p>{html.escape(piece)}</p>"
            for piece in content.split("\n\n")
            if piece.strip()
        )
        section_html_parts.append(f"<h2>{heading}</h2>{paragraph_html}")

    sections_html = "".join(section_html_parts)
    rendered_html = (
        '<!DOCTYPE html><html><head><meta charset="utf-8">'
        f"<title>{title}</title>"
        "<style>body{font-family:system-ui,sans-serif;max-width:760px;"
        "margin:2rem auto;padding:0 1rem;}h1{color:#4F46E5;}"
        "h2{color:#1f2937;margin-top:2rem;}p{line-height:1.6;}</style>"
        f"</head><body><h1>{title}</h1>{sections_html}</body></html>"
    )

    HTMLClass = _get_weasyprint_html()
    pdf_bytes: bytes = await asyncio.to_thread(
        HTMLClass(string=rendered_html).write_pdf,
    )
    return pdf_bytes


async def render_xlsx_from_source(source: dict[str, Any]) -> bytes:
    """Render a sheet-schema source dict to XLSX bytes.

    Args:
        source: Dict shaped ``{"sheets": [{"name": str, "rows": list[list[Any]]}]}``.
            Each sheet may also include optional ``"headers"`` and ``"title"``
            keys honored by :meth:`DocumentService._build_xlsx`.

    Returns:
        XLSX workbook bytes.

    Raises:
        ValueError: If ``source`` is missing the required ``"sheets"`` key.
    """
    if "sheets" not in source:
        msg = "render_xlsx_from_source: source is missing required key 'sheets'"
        raise ValueError(msg)

    return await asyncio.to_thread(
        DocumentService._build_xlsx,
        source["sheets"],
        {},
    )


async def render_pptx_from_source(source: dict[str, Any]) -> bytes:
    """Render a slide-JSON source dict to PPTX bytes.

    Args:
        source: Dict shaped ``{"title": str, "slides": [{"layout": str,
            "title": str, "body": str, "speaker_notes": str | None}]}``.
            ``body`` is split on newlines into bullet points; ``layout`` and
            ``speaker_notes`` are not yet wired into the rendered deck (deferred
            to v2 — see Section 7 of the document-viewer spec).

    Returns:
        PPTX presentation bytes.

    Raises:
        ValueError: If ``source`` is missing the required ``"slides"`` key.
    """
    if "slides" not in source:
        msg = "render_pptx_from_source: source is missing required key 'slides'"
        raise ValueError(msg)

    # Adapt the source shape (layout/title/body/speaker_notes) to the shape
    # expected by DocumentService._build_pptx (title/bullets/chart_image_bytes).
    # layout + speaker_notes are intentionally dropped here for v1.
    adapted_slides = [
        {
            "title": slide.get("title", ""),
            "bullets": [
                line
                for line in (slide.get("body") or "").split("\n")
                if line.strip()
            ],
        }
        for slide in source["slides"]
    ]

    return await asyncio.to_thread(
        DocumentService._build_pptx,
        adapted_slides,
        {},
    )


async def render_docx_from_source(source: dict[str, Any]) -> bytes:
    """Render a markdown-style source dict to DOCX bytes.

    Args:
        source: Dict shaped ``{"title": str, "sections": [{"heading": str,
            "content": str}]}``. ``content`` is split on blank lines into
            paragraphs.

    Returns:
        DOCX document bytes.

    Raises:
        ValueError: If ``source`` is missing the required ``"sections"`` key.
    """
    if "sections" not in source:
        msg = "render_docx_from_source: source is missing required key 'sections'"
        raise ValueError(msg)

    return await asyncio.to_thread(_build_docx_bytes, source)


def _build_docx_bytes(source: dict[str, Any]) -> bytes:
    """Build a DOCX document in memory (sync, runs in thread pool)."""
    from docx import Document

    doc = Document()
    doc.add_heading(str(source.get("title", "Untitled")), level=1)
    for section in source["sections"]:
        doc.add_heading(str(section.get("heading", "")), level=2)
        for paragraph in str(section.get("content", "")).split("\n\n"):
            if paragraph.strip():
                doc.add_paragraph(paragraph)
    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()
