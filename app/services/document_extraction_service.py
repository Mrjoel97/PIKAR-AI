# Copyright (c) 2024-2026 Pikar AI. All rights reserved.
# Proprietary and confidential. See LICENSE file for details.

"""Lazy extraction: binary -> text and binary -> canonical source.

``extract_text`` is called on the agent's first read of a user upload; it is
cheap and returns plain text. ``fork_to_source`` is called on the agent's first
edit; it is structured and produces the canonical source JSON used by the
re-render pipeline.

Per-format helpers:
- ``_pdf_to_text`` / ``_pdf_to_source`` -- pdfplumber
- ``_xlsx_to_text`` / ``_xlsx_to_source`` -- openpyxl (read_only + data_only
  for the text variant; data_only only for source so cached values, not
  formula text, are surfaced)
- ``_pptx_to_text`` / ``_pptx_to_source`` -- python-pptx
- ``_docx_to_text`` -- mammoth (raw text)
- ``_docx_to_source`` -- python-docx (sectioned by Heading style)

All CPU-bound work in the public ``extract_text`` / ``fork_to_source`` methods
is wrapped in ``asyncio.to_thread`` so we don't block the event loop.
"""

from __future__ import annotations

import asyncio
import io
from typing import Any

import mammoth
import pdfplumber
from openpyxl import load_workbook
from pptx import Presentation


class UnsupportedFormatError(ValueError):
    """Raised when a doc class is not handled by extraction.

    Subclasses :class:`ValueError` so callers catching ``ValueError`` still
    work without extra imports.
    """


class DocumentExtractionService:
    """Convert binary documents to plain text and to canonical source JSON.

    Stateless -- safe to share across requests. All public methods are
    coroutines; the actual decoding work runs in a worker thread via
    :func:`asyncio.to_thread` to keep the event loop responsive.
    """

    SUPPORTED_CLASSES: frozenset[str] = frozenset(
        {"report", "spreadsheet", "presentation", "word"}
    )

    async def extract_text(self, *, binary: bytes, doc_class: str) -> str:
        """Return plain text extracted from ``binary`` for ``doc_class``.

        Args:
            binary: The raw bytes of the source document.
            doc_class: One of ``report`` (PDF), ``spreadsheet`` (XLSX),
                ``presentation`` (PPTX), or ``word`` (DOCX).

        Returns:
            Plain-text representation of the document. May be empty for
            empty inputs (no exception raised).

        Raises:
            UnsupportedFormatError: If ``doc_class`` is not in
                :data:`SUPPORTED_CLASSES`.
        """
        if doc_class not in self.SUPPORTED_CLASSES:
            msg = f"Cannot extract text from doc_class={doc_class!r}"
            raise UnsupportedFormatError(msg)
        if doc_class == "report":
            return await asyncio.to_thread(self._pdf_to_text, binary)
        if doc_class == "spreadsheet":
            return await asyncio.to_thread(self._xlsx_to_text, binary)
        if doc_class == "presentation":
            return await asyncio.to_thread(self._pptx_to_text, binary)
        # doc_class == "word" -- the only remaining member of SUPPORTED_CLASSES
        return await asyncio.to_thread(self._docx_to_text, binary)

    async def fork_to_source(self, *, binary: bytes, doc_class: str) -> dict[str, Any]:
        """Return a canonical source dict extracted from ``binary``.

        Used by lazy fork on the agent's first edit: the binary is decoded
        into the same shape the renderer expects, so subsequent edits feed
        through the standard render pipeline.

        Args:
            binary: The raw bytes of the source document.
            doc_class: One of ``report`` (PDF), ``spreadsheet`` (XLSX),
                ``presentation`` (PPTX), or ``word`` (DOCX).

        Returns:
            A canonical source dict matching the v1 schema for the doc class.

        Raises:
            UnsupportedFormatError: If ``doc_class`` is not in
                :data:`SUPPORTED_CLASSES`.
        """
        if doc_class not in self.SUPPORTED_CLASSES:
            msg = f"Cannot fork doc_class={doc_class!r} to source"
            raise UnsupportedFormatError(msg)
        if doc_class == "report":
            return await asyncio.to_thread(self._pdf_to_source, binary)
        if doc_class == "spreadsheet":
            return await asyncio.to_thread(self._xlsx_to_source, binary)
        if doc_class == "presentation":
            return await asyncio.to_thread(self._pptx_to_source, binary)
        # doc_class == "word"
        return await asyncio.to_thread(self._docx_to_source, binary)

    # ------------------------------------------------------------------
    # Text helpers (sync; wrapped in to_thread by the public methods)
    # ------------------------------------------------------------------

    @staticmethod
    def _pdf_to_text(binary: bytes) -> str:
        """Extract concatenated page text from a PDF (pdfplumber)."""
        chunks: list[str] = []
        with pdfplumber.open(io.BytesIO(binary)) as pdf:
            for page in pdf.pages:
                text = page.extract_text() or ""
                chunks.append(text)
        return "\n\n".join(chunks)

    @staticmethod
    def _xlsx_to_text(binary: bytes) -> str:
        """Extract sheet/row text from an XLSX (read_only + data_only)."""
        wb = load_workbook(io.BytesIO(binary), read_only=True, data_only=True)
        chunks: list[str] = []
        for sheet in wb.worksheets:
            chunks.append(f"# {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                chunks.append("\t".join(str(c) if c is not None else "" for c in row))
        return "\n".join(chunks)

    @staticmethod
    def _pptx_to_text(binary: bytes) -> str:
        """Extract slide text frames in order, prefixed by slide headers."""
        prs = Presentation(io.BytesIO(binary))
        chunks: list[str] = []
        for i, slide in enumerate(prs.slides):
            chunks.append(f"--- Slide {i + 1} ---")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in paragraph.runs)
                        if text.strip():
                            chunks.append(text)
        return "\n".join(chunks)

    @staticmethod
    def _docx_to_text(binary: bytes) -> str:
        """Extract raw text from a DOCX via mammoth."""
        result = mammoth.extract_raw_text(io.BytesIO(binary))
        return result.value

    # ------------------------------------------------------------------
    # Source helpers (sync; wrapped in to_thread by the public methods)
    # ------------------------------------------------------------------

    @staticmethod
    def _pdf_to_source(binary: bytes) -> dict[str, Any]:
        """Convert a PDF to a report-shaped source (page -> section)."""
        # Naive: each page becomes a section with the whole page as content.
        # Better heading detection (font-size based) is deferred to follow-up.
        sections: list[dict[str, Any]] = []
        with pdfplumber.open(io.BytesIO(binary)) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                sections.append(
                    {
                        "heading": f"Page {i + 1}",
                        "content": text,
                    }
                )
        return {"title": "Imported PDF", "sections": sections}

    @staticmethod
    def _xlsx_to_source(binary: bytes) -> dict[str, Any]:
        """Convert an XLSX to a spreadsheet-shaped source (sheet -> rows)."""
        wb = load_workbook(io.BytesIO(binary), data_only=True)
        sheets: list[dict[str, Any]] = []
        for sheet in wb.worksheets:
            rows = [list(row) for row in sheet.iter_rows(values_only=True)]
            sheets.append({"name": sheet.title, "rows": rows})
        return {"sheets": sheets}

    @staticmethod
    def _pptx_to_source(binary: bytes) -> dict[str, Any]:
        """Convert a PPTX to a presentation-shaped source (slide -> dict)."""
        prs = Presentation(io.BytesIO(binary))
        slides: list[dict[str, Any]] = []
        for slide in prs.slides:
            title_text = ""
            body_text_parts: list[str] = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text
                placeholder_format = getattr(shape, "placeholder_format", None)
                if placeholder_format is not None and placeholder_format.idx == 0:
                    title_text = text
                else:
                    body_text_parts.append(text)
            speaker_notes: str | None = None
            if slide.has_notes_slide:
                speaker_notes = slide.notes_slide.notes_text_frame.text
            slides.append(
                {
                    "layout": "title_and_content",
                    "title": title_text,
                    "body": "\n".join(body_text_parts),
                    "speaker_notes": speaker_notes,
                }
            )
        return {"title": "Imported Deck", "slides": slides}

    @staticmethod
    def _docx_to_source(binary: bytes) -> dict[str, Any]:
        """Convert a DOCX to a word-shaped source (paragraphs grouped by Heading)."""
        from docx import Document  # inline to match document_service builders

        doc = Document(io.BytesIO(binary))
        sections: list[dict[str, Any]] = []
        current_heading = "Document"
        current_content: list[str] = []

        for para in doc.paragraphs:
            if para.style.name.startswith("Heading"):
                if current_content:
                    sections.append(
                        {
                            "heading": current_heading,
                            "content": "\n\n".join(current_content),
                        }
                    )
                    current_content = []
                current_heading = para.text
            elif para.text.strip():
                current_content.append(para.text)

        if current_content:
            sections.append(
                {
                    "heading": current_heading,
                    "content": "\n\n".join(current_content),
                }
            )

        return {"title": "Imported Document", "sections": sections}
